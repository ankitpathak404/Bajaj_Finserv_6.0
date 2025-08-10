import os
from pdfminer.high_level import extract_text
from docx import Document
import email
from email import policy
from email.parser import BytesParser
from qdrant_client import QdrantClient, models
from sentence_transformers import SentenceTransformer
import pandas as pd
from PIL import Image
import pytesseract
from pptx import Presentation
from io import BytesIO
import pdfplumber
# Set the tesseract executable path (adjust path as needed)
pytesseract.pytesseract.tesseract_cmd = r'C:/Program Files/Tesseract-OCR/tesseract.exe'

# Load the model

model = SentenceTransformer("BAAI/bge-large-en-v1.5",device="cuda")


client = QdrantClient(
    url="https://7ce5a2a4-bd12-4594-bf67-3add19a2a39b.us-west-2-0.aws.cloud.qdrant.io:6333", 
    api_key="eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJhY2Nlc3MiOiJtIn0.tVu-4QA6NR3Mv5AnRV-b5Rlz1QGz7tYqi9s2PRffLrA",
)


def extract_pdf_text(file_path):
    """
    Extracts structured text and tables from a PDF using pdfplumber
    """
    try:
        text_chunks = []
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text:
                    text_chunks.append(f"[Page {i+1} Text]:\n{text}")
                tables = page.extract_tables()
                for t_idx, table in enumerate(tables):
                    df = pd.DataFrame(table[1:], columns=table[0])
                    table_str = df.to_string(index=False)
                    text_chunks.append(f"[Page {i+1} Table {t_idx+1}]:\n{table_str}")
        return "\n\n".join(text_chunks)
    except Exception as e:
        print(f"[PDF ERROR] {file_path}: {e}")
        return ""


def extract_pptx_text(file_path):
    """
    Extracts text from a PPTX file including text from images using python-pptx and OCR
    """
    try:
        prs = Presentation(file_path)
        all_text = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = f"Slide {slide_num}:\n"
            
            for shape in slide.shapes:
                # Extract regular text
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
                
                # Extract text from images using OCR
                elif shape.shape_type == 13:  # Picture shape type
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        
                        # Convert bytes to PIL Image
                        from io import BytesIO
                        img = Image.open(BytesIO(image_bytes))
                        
                        # Extract text using OCR
                        ocr_text = pytesseract.image_to_string(img)
                        if ocr_text.strip():
                            slide_text += f"[Image Text]: {ocr_text}\n"
                    except Exception as img_error:
                        print(f"[PPTX IMAGE OCR ERROR] Slide {slide_num}: {img_error}")
                        continue
            
            all_text.append(slide_text)
        
        return "\n\n".join(all_text)
    except Exception as e:
        print(f"[PPTX ERROR] {file_path}: {e}")
        return ""

def extract_email_text(file_path):
    """
    Extracts text from a .eml email file using Python's email parser
    """
    try:
        with open(file_path, 'rb') as f:
            msg = BytesParser(policy=policy.default).parse(f)
        body = msg.get_body(preferencelist=('plain'))
        return body.get_content() if body else ""
    except Exception as e:
        print(f"[EMAIL ERROR] {file_path}: {e}")
        return ""

def extract_xlsx_text(file_path):
    try:
        xl_file = pd.ExcelFile(file_path)
        all_text = []
        for sheet_name in xl_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            sheet_text = f"[Sheet: {sheet_name}]\n{df.to_string(index=False)}"
            all_text.append(sheet_text)
        return "\n\n".join(all_text)
    except Exception as e:
        print(f"[XLSX ERROR] {file_path}: {e}")
        return ""


def extract_image_text(file_path):
    """
    Extracts text from image files (JPEG, PNG) using OCR (pytesseract)
    """
    try:
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        print(f"[IMAGE ERROR] {file_path}: {e}")
        return ""
    
def extract_docx_text(file_path):
    try:
        doc = Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            style = para.style.name.lower()
            if 'heading' in style:
                full_text.append(f"[Heading]: {para.text}")
            elif 'list' in style:
                full_text.append(f"[List Item]: {para.text}")
            elif para.text.strip():
                full_text.append(para.text)

        # Table extraction
        for table_idx, table in enumerate(doc.tables):
            data = []
            for row in table.rows:
                data.append([cell.text.strip() for cell in row.cells])
            df = pd.DataFrame(data[1:], columns=data[0]) if len(data) > 1 else pd.DataFrame(data)
            full_text.append(f"[Table {table_idx+1}]:\n{df.to_string(index=False)}")

        return "\n\n".join(full_text)
    except Exception as e:
        print(f"[DOCX ERROR] {file_path}: {e}")
        return ""


#payload={"document_name":"extracted_text"}
def extract_text_from_file(file_path):
    """
    General handler based on file type
    """
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_pdf_text(file_path)
    elif ext == ".docx":
        return extract_docx_text(file_path)
    elif ext == ".pptx":
        return extract_pptx_text(file_path)
    elif ext == ".eml":
        return extract_email_text(file_path)
    elif ext == ".xlsx":
        return extract_xlsx_text(file_path)
    elif ext in [".jpeg", ".jpg", ".png", ".bmp", ".tiff", ".tif"]:
        return extract_image_text(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

from langchain.text_splitter import RecursiveCharacterTextSplitter

def langchain_chunk(text):
    splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
        chunk_size=600,           # adjust as needed (1000 for vector DBs)
        chunk_overlap=150,
        model_name="gpt-3.5-turbo",
        separators=["\n\n", "\n", ".", " ", ""]
    )
    return splitter.split_text(text)


if __name__ == "__main__":
    # Example usage
    test_files = [
            
        
    ]
    
    idx=0

    for file in test_files:
        if not os.path.exists(file):
            print(f"[SKIP] File not found: {file}")
            continue

        print(f"\n=== Extracting: {file} ===")
        text = extract_text_from_file(file)
        text=langchain_chunk(text)
        
        for i in range(len(text)): 

            print(idx,"\n")
            v=model.encode(text[i])
            client.upsert(
                collection_name="bajaj2",
                points=[
                 models.PointStruct(
            id=idx,
            payload={
                str(file): str(text[i]),
                "file_name":str(file),
                "chunk_id": i,
                "length": len(text[i]),
                "chunk_preview": text[i][:100]     
                     
                     },
            vector=v,  # ← just a list, not a dict
            )
            ])
            idx+=1

        text = "\n\n".join(text)  # Join chunks with double newlines for better readability
        # Save to .txt file
        txt_filename = os.path.splitext(file)[0] + ".txt"
        with open(txt_filename, "w", encoding="utf-8") as f:
            f.write(text)

        print(f"[SAVED] Extracted text saved to: {txt_filename}")
        print("\n--- End of Output ---\n")